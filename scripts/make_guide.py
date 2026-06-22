# -*- coding: utf-8 -*-
"""IShop ашиглах заавар — PowerPoint гарын авлага үүсгэгч (screenshot дэмжсэн)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ── Өнгөний загвар (dark zinc / teal) ──────────────────────────────────────
BG       = RGBColor(0x0A, 0x0A, 0x0B)
CARD     = RGBColor(0x18, 0x18, 0x1B)
CARD2    = RGBColor(0x27, 0x27, 0x2A)
TEAL     = RGBColor(0x14, 0xB8, 0xA6)
TEAL_LT  = RGBColor(0x5E, 0xEA, 0xD4)
WHITE    = RGBColor(0xFA, 0xFA, 0xFA)
GREY     = RGBColor(0xA1, 0xA1, 0xAA)
GREY_DK  = RGBColor(0x71, 0x71, 0x7A)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
VIOLET   = RGBColor(0x8B, 0x5C, 0xF6)
BLUE     = RGBColor(0x3B, 0x82, 0xF6)
RED      = RGBColor(0xEF, 0x44, 0x44)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)

FONT = "Segoe UI"
SHOT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "screenshots")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

_NUM = 0
def nextnum():
    global _NUM
    _NUM += 1
    return _NUM


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.10):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = FONT
    return tb


def chip(s, x, y, label, color=TEAL):
    w = 0.34 + 0.115 * len(label)
    box(s, x, y, w, 0.42, fill=None, line=color, line_w=1.25, radius=0.5)
    txt(s, x, y - 0.02, w, 0.46, [[(label, 12, color, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    return w


def header(s, kicker, title, num):
    box(s, 0.0, 0.0, 0.16, 7.5, fill=TEAL, radius=0)
    txt(s, 0.75, 0.55, 11, 0.4, [[(kicker, 13, TEAL, True)]], space_after=0)
    txt(s, 0.72, 0.92, 11.5, 1.0, [[(title, 32, WHITE, True)]], space_after=0)
    txt(s, 12.2, 0.6, 0.9, 0.5, [[(f"{num:02d}", 15, GREY_DK, True)]],
        align=PP_ALIGN.RIGHT, space_after=0)


def feature_card(s, x, y, w, h, icon, title, lines, accent=TEAL):
    box(s, x, y, w, h, fill=CARD, line=CARD2, line_w=1.0, radius=0.09)
    txt(s, x + 0.3, y + 0.25, 1.0, 0.6, [[(icon, 26, accent, False)]], space_after=0)
    txt(s, x + 0.3, y + 0.95, w - 0.6, 0.5, [[(title, 16, WHITE, True)]], space_after=0)
    body = [[(ln, 11.5, GREY, False)] for ln in lines]
    txt(s, x + 0.3, y + 1.45, w - 0.55, h - 1.6, body, space_after=5, line_spacing=1.05)


def step_card(s, x, y, w, n, title, desc, accent=TEAL):
    h = 1.18
    box(s, x, y, w, h, fill=CARD, line=CARD2, line_w=1.0, radius=0.12)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.28), Inches(y + 0.34), Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = accent
    circ.line.fill.background(); circ.shadow.inherit = False
    pc = circ.text_frame.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
    rc = pc.add_run(); rc.text = str(n)
    rc.font.size = Pt(18); rc.font.bold = True; rc.font.color.rgb = BG; rc.font.name = FONT
    txt(s, x + 1.0, y + 0.18, w - 1.2, 0.5, [[(title, 15, WHITE, True)]], space_after=0)
    txt(s, x + 1.0, y + 0.58, w - 1.2, h - 0.6, [[(desc, 11.5, GREY, False)]],
        space_after=0, line_spacing=1.05)


# ── Screenshot frame helpers ────────────────────────────────────────────────
def browser_frame(s, x, y, w, h):
    """Browser-style framed card; returns inner image area (x,y,w,h) in inches."""
    box(s, x, y, w, h, fill=CARD, line=CARD2, line_w=1.0, radius=0.04)
    bar_h = 0.36
    box(s, x, y, w, bar_h, fill=CARD2, radius=0.04)
    for i, c in enumerate([RED, AMBER, GREEN]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.22 + i * 0.2),
                               Inches(y + 0.12), Inches(0.11), Inches(0.11))
        d.fill.solid(); d.fill.fore_color.rgb = c
        d.line.fill.background(); d.shadow.inherit = False
    pad = 0.1
    return x + pad, y + bar_h + pad, w - 2 * pad, h - bar_h - 2 * pad


def fit_image(s, path, bx, by, bw, bh):
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = bw / bh
    if ar > box_ar:
        w = bw; h = bw / ar
    else:
        h = bh; w = bh * ar
    x = bx + (bw - w) / 2
    y = by + (bh - h) / 2
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))


def placeholder(s, bx, by, bw, bh, fname):
    box(s, bx, by, bw, bh, fill=BG, line=TEAL, line_w=1.5, radius=0.03)
    txt(s, bx, by + bh / 2 - 0.75, bw, 1.5,
        [[("🖼", 34, TEAL, False)],
         [("Энд зураг орно", 14, GREY, False)],
         [(fname, 13, TEAL_LT, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=5)


def shot_slide(kicker, title, shots):
    """shots: list of (filename, caption). 1 → том, 2 → зэрэгцээ."""
    s = slide()
    header(s, kicker, title, nextnum())
    if len(shots) == 1:
        fname, cap = shots[0]
        fx, fy, fw, fh = 1.85, 1.95, 9.6, 4.5
        ix, iy, iw, ih = browser_frame(s, fx, fy, fw, fh)
        p = os.path.join(SHOT_DIR, fname)
        if os.path.exists(p):
            fit_image(s, p, ix, iy, iw, ih)
        else:
            placeholder(s, ix, iy, iw, ih, fname)
        txt(s, fx, fy + fh + 0.12, fw, 0.5, [[("› " + cap, 13.5, GREY, False)]],
            align=PP_ALIGN.CENTER, space_after=0)
    else:
        positions = [(0.75, 5.7), (7.0, 5.7)]
        for (fname, cap), (fx, w) in zip(shots, positions):
            fy, fh = 1.95, 4.3
            ix, iy, iw, ih = browser_frame(s, fx, fy, w, fh)
            p = os.path.join(SHOT_DIR, fname)
            if os.path.exists(p):
                fit_image(s, p, ix, iy, iw, ih)
            else:
                placeholder(s, ix, iy, iw, ih, fname)
            txt(s, fx, fy + fh + 0.12, w, 0.5, [[("› " + cap, 12.5, GREY, False)]],
                align=PP_ALIGN.CENTER, space_after=0)
    return s


# ════════════════════════════════════════════════════════════════════════════
# 1. ГАРЧИГ
# ════════════════════════════════════════════════════════════════════════════
s = slide(); nextnum()
box(s, 9.3, -1.2, 5.5, 5.5, fill=CARD, radius=0.5)
box(s, 10.5, 4.3, 4.5, 4.5, fill=CARD, radius=0.5)
txt(s, 0.9, 1.5, 8, 0.5, [[("ХЭРЭГЛЭГЧИЙН ГАРЫН АВЛАГА", 15, TEAL, True)]], space_after=0)
txt(s, 0.85, 2.05, 10, 1.6, [[("ISHOP ", 60, WHITE, True), ("платформ", 60, TEAL, True)]], space_after=0)
txt(s, 0.9, 3.5, 9.5, 1.4,
    [[("Онлайн худалдааны систем ашиглах бүрэн заавар —", 18, GREY, False)],
     [("хэрэглэгч болон админ хэсгийн дэлгэрэнгүй гарын авлага", 18, GREY, False)]],
    space_after=4, line_spacing=1.1)
w1 = chip(s, 0.9, 5.4, "Next.js 16")
chip(s, 0.9 + w1 + 0.2, 5.4, "Хэрэглэгчийн заавар", TEAL_LT)
txt(s, 0.9, 6.6, 11, 0.5, [[("IShop · Онлайн дэлгүүр · ", 12, GREY_DK, False), ("2026", 12, TEAL, True)]], space_after=0)

# ════════════════════════════════════════════════════════════════════════════
# 2. АГУУЛГА
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "ЕРӨНХИЙ", "Энэ гарын авлагад юу багтсан бэ?", nextnum())
items = [
    ("01", "Системийн танилцуулга", "IShop гэж юу вэ, хэн ашиглах вэ"),
    ("02", "Бүртгэл ба нэвтрэлт", "Шинээр бүртгүүлэх, нэвтрэх, нууц үг сэргээх"),
    ("03", "Дэлгүүр хэсэг (хэрэглэгч)", "Бараа хайх, үзэх, сагсалж худалдан авах"),
    ("04", "Захиалга ба профайл", "Захиалга хянах, хаяг, хүслийн жагсаалт"),
    ("05", "Админ — Хянах самбар", "Орлого, статистик, захиалгын төлөв"),
    ("06", "Админ — Бүтээгдэхүүн", "Бараа нэмэх, засах, ангилал, онцлох"),
    ("07", "Админ — Захиалга & Санхүү", "Захиалга боловсруулах, төлбөр, тайлан"),
    ("08", "Зөвлөмж ба тусламж", "Түгээмэл асуулт, анхаарах зүйлс"),
]
x0, y0, cw, ch, gx, gy = 0.75, 1.95, 5.85, 1.05, 0.3, 0.18
for i, (n, t, d) in enumerate(items):
    col = i % 2; row = i // 2
    x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
    box(s, x, y, cw, ch, fill=CARD, line=CARD2, line_w=1.0, radius=0.12)
    txt(s, x + 0.28, y + 0.16, 0.9, 0.7, [[(n, 22, TEAL, True)]], space_after=0)
    txt(s, x + 1.15, y + 0.16, cw - 1.3, 0.4, [[(t, 14.5, WHITE, True)]], space_after=0)
    txt(s, x + 1.15, y + 0.58, cw - 1.3, 0.4, [[(d, 11, GREY, False)]], space_after=0)

# ════════════════════════════════════════════════════════════════════════════
# 3. ТАНИЛЦУУЛГА
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "01 · ТАНИЛЦУУЛГА", "IShop гэж юу вэ?", nextnum())
txt(s, 0.75, 1.95, 11.8, 1.0,
    [[("IShop нь онлайн худалдааны бүрэн систем юм. Хэрэглэгч бараа сонгож худалдан авдаг ", 15, GREY, False),
      ("дэлгүүр", 15, TEAL, True),
      (" хэсэг, дэлгүүрийг удирддаг ", 15, GREY, False),
      ("админ", 15, TEAL, True),
      (" хэсэг гэсэн хоёр үндсэн талтай.", 15, GREY, False)]],
    space_after=0, line_spacing=1.15)
feature_card(s, 0.75, 3.1, 5.85, 3.7, "🛍️", "Хэрэглэгчийн хэсэг (дэлгүүр)",
             ["• Бараа хайх, ангиллаар шүүх", "• Дэлгэрэнгүй мэдээлэл, зураг үзэх",
              "• Сагсанд нэмж худалдан авах", "• QPay-ээр төлбөр төлөх",
              "• Захиалгаа хянах, хаяг хадгалах", "• Хүслийн жагсаалт үүсгэх"], accent=TEAL)
feature_card(s, 6.95, 3.1, 5.6, 3.7, "⚙️", "Админ хэсэг (удирдлага)",
             ["• Хянах самбар, статистик", "• Бараа, ангилал нэмэх/засах",
              "• Захиалга боловсруулах", "• Харилцагчдыг удирдах",
              "• Төлбөр, нэхэмжлэл, тайлан", "• Мэдэгдэл, тохиргоо"], accent=VIOLET)

# ════════════════════════════════════════════════════════════════════════════
# 4. БҮРТГЭЛ БА НЭВТРЭЛТ  (+ screenshot)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "02 · ЭХЛЭЛ", "Бүртгүүлэх ба нэвтрэх", nextnum())
step_card(s, 0.75, 1.95, 5.85, 1, "Бүртгүүлэх", "Нүүр хуудаснаас «Бүртгүүлэх» сонгож, нэр, и-мэйл, нууц үгээ оруулна.")
step_card(s, 0.75, 3.31, 5.85, 2, "Нэвтрэх", "И-мэйл, нууц үгээрээ нэвтэрнэ. Амжилттай бол шууд дэлгүүр рүү орно.")
step_card(s, 0.75, 4.67, 5.85, 3, "Нууц үг сэргээх", "Мартсан бол «Нууц үг сэргээх»-ээр и-мэйлээр код авч шинэчилнэ.")
box(s, 6.95, 1.95, 5.6, 4.85, fill=CARD, line=CARD2, line_w=1.0, radius=0.07)
txt(s, 7.3, 2.25, 5, 0.5, [[("💡 Анхаарах", 16, TEAL, True)]], space_after=0)
txt(s, 7.3, 2.85, 4.95, 3.8,
    [[("• И-мэйл хаяг зөв, идэвхтэй байх ёстой — баталгаажуулалт, мэдэгдэл ирнэ.", 12.5, GREY, False)],
     [("• Нууц үгээ найдвартай (8+ тэмдэгт) сонгоно.", 12.5, GREY, False)],
     [("• Нэг бүртгэлээр олон төхөөрөмжөөс нэвтэрч болно («Төхөөрөмжүүд» хэсэг).", 12.5, GREY, False)],
     [("• Админ эрхтэй бол нэвтрэхэд автоматаар админ хэсэг рүү шилжинэ.", 12.5, GREY, False)]],
    space_after=10, line_spacing=1.1)
shot_slide("02 · ЭХЛЭЛ", "Нэвтрэх / бүртгүүлэх дэлгэц", [("login.png", "Нэвтрэх ба бүртгүүлэх хуудас")])

# ════════════════════════════════════════════════════════════════════════════
# 5. ДЭЛГҮҮР — ХАЙХ, ҮЗЭХ  (+ screenshots)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "03 · ХЭРЭГЛЭГЧ", "Бараа хайх ба үзэх", nextnum())
feature_card(s, 0.75, 1.95, 3.75, 4.85, "🏠", "Нүүр хуудас",
             ["Онцлох бараа, шилдэг", "ангилал, статистик", "харагдана.",
              "", "Дээд хэсгийн хайлтын", "талбараас нэр, тайлбар,", "ангиллаар хайна."])
feature_card(s, 4.75, 1.95, 3.75, 4.85, "🔍", "Хайх & шүүх",
             ["Ангиллаар шүүж, нэрээр", "хайж барааг хурдан", "олно.",
              "", "Үр дүн шууд шинэчлэгдэж", "тохирох бараанууд", "жагсаалтаар гарна."], accent=BLUE)
feature_card(s, 8.75, 1.95, 3.8, 4.85, "📄", "Дэлгэрэнгүй",
             ["Барааны зураг, үнэ,", "тайлбар, нөөцийг", "харна.",
              "", "Тоо ширхэг сонгож", "«Сагсанд нэмэх» эсвэл", "хүслийн жагсаалтад нэмнэ."], accent=AMBER)
shot_slide("03 · ХЭРЭГЛЭГЧ", "Дэлгүүрийн дэлгэцүүд",
           [("user_products.png", "Дэлгүүр — барааны жагсаалт, хайлт"),
            ("product_details.png", "Барааны дэлгэрэнгүй хуудас")])

# ── ХУВИЛБАР СОНГОХ (шинэ дэлгэрэнгүй) ───────────────────────────────────────
s = slide()
header(s, "03 · ХЭРЭГЛЭГЧ", "Хувилбартай бараа сонгох", nextnum())
feature_card(s, 0.75, 1.95, 3.75, 4.85, "🎨", "Өнгө / хэмжээ",
             ["Нэг бараа олон өнгө,", "хэмжээ, материал,", "загвартай байж болно.",
              "", "Товчлуураар хүссэн", "хувилбараа сонгоно."], accent=TEAL)
feature_card(s, 4.75, 1.95, 3.75, 4.85, "📦", "Бодит нөөц & үнэ",
             ["Сонгосон хослолын нөөц,", "үнэ шууд шинэчлэгдэнэ.",
              "", "Боломжгүй хослолыг", "автоматаар тохирох", "хувилбар руу тааруулна."], accent=BLUE)
feature_card(s, 8.75, 1.95, 3.8, 4.85, "✨", "Онцлог шинж",
             ["WiFi, Bluetooth, Type-C", "зэрэг нэмэлт шинжийг", "доор харуулна.",
              "", "Сонгоод «Сагсанд", "нэмэх» дарна."], accent=AMBER)
shot_slide("03 · ХЭРЭГЛЭГЧ", "Барааны дэлгэрэнгүй — хувилбар сонгох",
           [("product_details_new_user.png", "Өнгө, хэмжээ, материал сонгоход нөөц, үнэ шинэчлэгдэнэ")])

# ════════════════════════════════════════════════════════════════════════════
# 6. ДЭЛГҮҮР — САГС, ХУДАЛДАН АВАЛТ  (+ screenshots)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "03 · ХЭРЭГЛЭГЧ", "Сагс ба худалдан авалт", nextnum())
steps = [
    ("Сагсанд нэмэх", "Барааг тоо ширхэгтэйгээр сагсанд нэмнэ."),
    ("Сагс шалгах", "Сагсан дахь бараа, тоо, нийт дүнг хянана."),
    ("Хаяг сонгох", "Хүргэлтийн хаягаа оруулах буюу сонгоно."),
    ("Төлбөр төлөх", "QPay-ээр төлбөрөө төлж захиалгаа баталгаажуулна."),
]
for i, (t, d) in enumerate(steps):
    step_card(s, 0.75, 1.95 + i * 1.25, 7.0, i + 1, t, d)
box(s, 8.1, 1.95, 4.45, 4.85, fill=CARD, line=TEAL, line_w=1.25, radius=0.07)
txt(s, 8.45, 2.25, 3.8, 0.5, [[("🧾 Төлбөрийн урсгал", 15, TEAL, True)]], space_after=0)
txt(s, 8.45, 2.95, 3.8, 3.7,
    [[("Худалдан авалт 4 алхамтай: сагс → хаяг → төлбөр → баталгаажуулалт.", 12.5, GREY, False)],
     [("QPay-ийн QR кодыг банкны аппаар уншуулж төлнө.", 12.5, GREY, False)],
     [("Төлбөр амжилттай бол захиалга «Хүлээгдэж буй» төлөвт орж, админд мэдэгдэнэ.", 12.5, GREY, False)]],
    space_after=12, line_spacing=1.12)
shot_slide("03 · ХЭРЭГЛЭГЧ", "Сагс ба хаяг сонгох",
           [("cart.png", "Сагс — бараа, тоо, нийт дүн"),
            ("address_pick_step.png", "Хүргэлтийн хаяг сонгох алхам")])
shot_slide("03 · ХЭРЭГЛЭГЧ", "Төлбөр төлөх алхам",
           [("payment_method_step.png", "Төлбөрийн хэлбэр сонгох"),
            ("qpay_step.png", "QPay QR кодоор төлөх")])

# ════════════════════════════════════════════════════════════════════════════
# 7. ЗАХИАЛГА БА ПРОФАЙЛ  (+ screenshot)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "04 · ХЭРЭГЛЭГЧ", "Захиалга & хувийн мэдээлэл", nextnum())
cards = [
    ("📦", "Захиалгууд", ["Захиалгын түүх,", "төлөв, дэлгэрэнгүйг", "хянана."], TEAL),
    ("❤️", "Хүслийн жагсаалт", ["Дараа авах барааг", "хадгалж, шууд сагсанд", "шилжүүлнэ."], RGBColor(0xF4, 0x3F, 0x5E)),
    ("📍", "Хаяг", ["Хүргэлтийн хаягаа", "хадгалж, газрын зураг", "дээр тэмдэглэнэ."], BLUE),
    ("🔔", "Мэдэгдэл", ["Захиалга, урамшууллын", "мэдэгдлийг хүлээн", "авна."], AMBER),
    ("👤", "Профайл", ["Нэр, зураг, мэдээллээ", "засварлана."], VIOLET),
    ("⚙️", "Тохиргоо", ["Нууц үг, төхөөрөмж,", "хэл, харагдацыг", "тохируулна."], TEAL_LT),
]
x0, y0, cw, ch, gx, gy = 0.75, 1.95, 3.75, 2.3, 0.18, 0.2
for i, (ic, t, lines, ac) in enumerate(cards):
    feature_card(s, x0 + (i % 3) * (cw + gx), y0 + (i // 3) * (ch + gy), cw, ch, ic, t, lines, accent=ac)
shot_slide("04 · ХЭРЭГЛЭГЧ", "Захиалгын түүх ба дэлгэрэнгүй",
           [("user_orders.png", "Миний захиалгууд"),
            ("user_order_detail.png", "Захиалгын дэлгэрэнгүй")])

# ── САЛБАРУУД (хэрэглэгч) ─────────────────────────────────────────────────────
s = slide()
header(s, "04 · ХЭРЭГЛЭГЧ", "Салбар, байршлууд", nextnum())
txt(s, 0.75, 1.9, 11.8, 0.7,
    [[("Дэлгүүрийн салбаруудыг ", 14, GREY, False), ("«Салбарууд»", 14, TEAL, True),
      (" хуудаснаас үзнэ. Дээд талын товчоор ", 14, GREY, False),
      ("жагсаалт", 14, TEAL, True), (" эсвэл ", 14, GREY, False),
      ("газрын зураг", 14, TEAL, True), (" хэлбэрээр сэлгэнэ.", 14, GREY, False)]],
    space_after=0, line_spacing=1.12)
feature_card(s, 0.75, 2.95, 5.85, 3.85, "📋", "Жагсаалт хэлбэр",
             ["Салбар бүрийн нэр, хаяг,", "утсыг карт хэлбэрээр", "харна.",
              "", "«Газрын зураг дээр харах»-аар", "тухайн салбарын байршлыг", "шууд нээнэ."], accent=TEAL)
feature_card(s, 6.95, 2.95, 5.6, 3.85, "🗺️", "Газрын зураг хэлбэр",
             ["Бүх салбарыг Google", "газрын зураг дээр пин", "болгон харуулна.",
              "", "Пин дээр дарж нэр, хаяг,", "утсыг шууд харна."], accent=BLUE)
shot_slide("04 · ХЭРЭГЛЭГЧ", "Салбарууд — жагсаалт ба газрын зураг",
           [("user_see_branches_card_mode.png", "Жагсаалт хэлбэр — салбарын картууд"),
            ("user_see_branches_map_mode.png", "Газрын зураг хэлбэр — Google Maps пин")])

# ════════════════════════════════════════════════════════════════════════════
# 8. АДМИН — ХЯНАХ САМБАР  (+ screenshot)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "05 · АДМИН", "Хянах самбар (Dashboard)", nextnum())
txt(s, 0.75, 1.9, 11.8, 0.6,
    [[("Админ эрхээр нэвтрэхэд ", 14, GREY, False), ("ISHOP ADMIN", 14, TEAL, True),
      (" удирдлагын самбар нээгдэнэ. Зүүн талын цэснээс бүх хэсэгт шилжинэ.", 14, GREY, False)]],
    space_after=0, line_spacing=1.1)
stats = [("💰", "Нийт орлого", TEAL), ("📦", "Нийт захиалга", BLUE),
         ("👥", "Хэрэглэгчид", VIOLET), ("🏬", "Нийт бараа", AMBER)]
for i, (ic, t, ac) in enumerate(stats):
    x = 0.75 + i * 3.02
    box(s, x, 2.75, 2.82, 1.5, fill=CARD, line=CARD2, line_w=1.0, radius=0.14)
    txt(s, x + 0.28, 2.95, 1, 0.6, [[(ic, 24, ac, False)]], space_after=0)
    txt(s, x + 0.28, 3.65, 2.4, 0.5, [[(t, 13, WHITE, True)]], space_after=0)
box(s, 0.75, 4.5, 7.4, 2.3, fill=CARD, line=CARD2, line_w=1.0, radius=0.08)
txt(s, 1.1, 4.75, 6.6, 0.5, [[("📊 Самбарт харагдах зүйлс", 15, TEAL, True)]], space_after=0)
txt(s, 1.1, 5.35, 6.7, 1.4,
    [[("• Орлогын график (огноогоор шүүж болно)", 12.5, GREY, False)],
     [("• Захиалгын төлөв (хүлээгдэж буй / хүргэгдсэн)", 12.5, GREY, False)],
     [("• Шилдэг борлуулалттай бараа ба сүүлийн захиалгууд", 12.5, GREY, False)]],
    space_after=8, line_spacing=1.1)
box(s, 8.4, 4.5, 4.15, 2.3, fill=CARD, line=AMBER, line_w=1.25, radius=0.08)
txt(s, 8.75, 4.75, 3.5, 0.5, [[("⏳ Хурдан анхаарал", 15, AMBER, True)]], space_after=0)
txt(s, 8.75, 5.35, 3.5, 1.4,
    [[("Хүлээгдэж буй захиалга байвал шар сэрэмжлүүлэг гарч, дарахад шууд «Захиалгууд» руу шилжинэ.", 12.5, GREY, False)]],
    space_after=0, line_spacing=1.12)
shot_slide("05 · АДМИН", "Админ хянах самбар", [("admin_dashboard.png", "Хянах самбар — статистик ба график")])

# ════════════════════════════════════════════════════════════════════════════
# 9. АДМИН — БҮТЭЭГДЭХҮҮН  (+ screenshot)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "06 · АДМИН", "Бүтээгдэхүүний удирдлага", nextnum())
cards = [
    ("📋", "Жагсаалт", ["Бүх барааг харах,", "хайх, засах, устгах."], TEAL),
    ("➕", "Шинэ нэмэх", ["Нэр, үнэ, нөөц,", "зураг, ангилал оруулж", "шинэ бараа үүсгэх."], BLUE),
    ("⭐", "Онцлох", ["Нүүр хуудсанд гарах", "онцлох барааг сонгох."], AMBER),
    ("🗂️", "Ангилал", ["Барааны ангилал нэмэх,", "зураг тохируулах,", "засварлах."], VIOLET),
]
x0, cw, gx = 0.75, 2.85, 0.13
for i, (ic, t, lines, ac) in enumerate(cards):
    feature_card(s, x0 + i * (cw + gx), 2.05, cw, 3.4, ic, t, lines, accent=ac)
box(s, 0.75, 5.7, 11.8, 1.1, fill=CARD, line=TEAL, line_w=1.25, radius=0.1)
txt(s, 1.1, 5.92, 11.2, 0.7,
    [[("💡 Зөвлөмж:  ", 13, TEAL, True),
      ("Шинэ бараа нэмэхдээ чанартай зураг, тодорхой тайлбар, зөв ангилал, бодит нөөцийн тоог оруулаарай — энэ нь хэрэглэгчид олж, итгэхэд тусална.", 12.5, GREY, False)]],
    space_after=0, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
shot_slide("06 · АДМИН", "Бүтээгдэхүүн ба ангилал",
           [("admin_products.png", "Барааны жагсаалт"),
            ("admin_categories.png", "Ангиллын удирдлага")])

# ── ХУВИЛБАРТАЙ БАРАА ҮҮСГЭХ (админ) ─────────────────────────────────────────
s = slide()
header(s, "06 · АДМИН", "Хувилбартай бараа (өнгө, хэмжээ)", nextnum())
step_card(s, 0.75, 1.95, 11.8, 1, "Шинж чанар тодорхойлох",
          "Өнгө, хэмжээ, загвар, материал бүрийн утгуудыг нэмнэ (ж: Улаан, XL, Cotton).")
step_card(s, 0.75, 3.31, 11.8, 2, "Хувилбар (хослол) үүсгэх",
          "Утгуудыг хослуулан хувилбар нэмж, тус бүрийн нөөц, үнэ, хямдрал, SKU-г оруулна.")
step_card(s, 0.75, 4.67, 11.8, 3, "Онцлог шинж нэмэх",
          "WiFi, Bluetooth, Type-C зэрэг нэмэлт шинжийг жагсаалтаар оруулна.")
box(s, 0.75, 6.05, 11.8, 0.8, fill=CARD, line=TEAL, line_w=1.25, radius=0.13)
txt(s, 1.1, 6.18, 11.2, 0.6,
    [[("💡 ", 12.5, TEAL, True),
      ("Хувилбарын үнэ хоосон бол барааны нийт үнэ хэрэглэгдэнэ. Нөөц 0 бол тухайн хослол хэрэглэгчид сонгогдохгүй.", 12.5, GREY, False)]],
    space_after=0, anchor=MSO_ANCHOR.MIDDLE)
shot_slide("06 · АДМИН", "Хувилбар тодорхойлох ба нөөц/үнэ",
           [("product_variants.png", "Шинж чанар, утгууд тодорхойлох"),
            ("product_variant_and_stock.png", "Хувилбар бүрийн нөөц, үнэ, онцлог шинж")])

# ════════════════════════════════════════════════════════════════════════════
# 10. АДМИН — ЗАХИАЛГА & САНХҮҮ  (+ screenshots)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "07 · АДМИН", "Захиалга & санхүү", nextnum())
feature_card(s, 0.75, 1.95, 3.75, 4.85, "🛒", "Захиалгууд",
             ["Ирсэн захиалгыг харах,", "дэлгэрэнгүй нээх,", "төлөв шинэчлэх",
              "(боловсруулах →", "хүргэх → дуусгах).", "", "Хэрэглэгчид автоматаар", "мэдэгдэл очно."], accent=TEAL)
feature_card(s, 4.75, 1.95, 3.75, 4.85, "💳", "Төлбөр & Invoice",
             ["«Төлбөрүүд» — бүх", "гүйлгээг хянана.", "",
              "«QPay Invoice» —", "нэхэмжлэлийн төлөв,", "QR кодыг шалгана."], accent=BLUE)
feature_card(s, 8.75, 1.95, 3.8, 4.85, "📈", "Тайлан",
             ["Борлуулалт, орлогын", "тайланг хугацаагаар", "гаргана.",
              "", "Excel рүү экспортлож", "хадгалах боломжтой", "(.xlsx)."], accent=AMBER)
shot_slide("07 · АДМИН", "Захиалга ба тайлан",
           [("admin_orders.png", "Захиалгын удирдлага"),
            ("admin_report.png", "Санхүүгийн тайлан")])

# ════════════════════════════════════════════════════════════════════════════
# 11. АДМИН — ХАРИЛЦАГЧ, ТОХИРГОО, МЭДЭГДЭЛ  (+ screenshot)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "07 · АДМИН", "Харилцагч · Тохиргоо · Мэдэгдэл", nextnum())
cards = [
    ("👥", "Харилцагчид", ["Бүртгэлтэй хэрэглэгчдийн", "жагсаалт, захиалгын", "түүх, дэлгэрэнгүйг үзэх."], VIOLET),
    ("🔔", "Мэдэгдэл", ["Шинэ захиалга, үйл", "явдлын мэдэгдлийг", "хүлээн авах, хянах."], AMBER),
    ("⚙️", "Тохиргоо", ["Дэлгүүрийн мэдээлэл,", "админ профайл, системийн", "тохиргоог удирдах."], TEAL),
]
x0, cw, gx = 0.75, 3.85, 0.13
for i, (ic, t, lines, ac) in enumerate(cards):
    feature_card(s, x0 + i * (cw + gx), 2.05, cw, 3.6, ic, t, lines, accent=ac)
box(s, 0.75, 5.95, 11.8, 0.85, fill=CARD, line=CARD2, line_w=1.0, radius=0.13)
txt(s, 1.1, 6.1, 11.2, 0.6,
    [[("Цэс:  ", 12.5, GREY_DK, True),
      ("Хянах самбар · Захиалгууд · Бүтээгдэхүүн · Харилцагчид · Санхүү · Тохиргоо · Мэдэгдэл", 12.5, GREY, False)]],
    space_after=0, anchor=MSO_ANCHOR.MIDDLE)
shot_slide("07 · АДМИН", "Харилцагчид ба тохиргоо",
           [("admin_users.png", "Харилцагчдын жагсаалт"),
            ("admin_settings.png", "Системийн тохиргоо")])

# ── САЛБАРЫН ТОХИРГОО (админ) ─────────────────────────────────────────────────
s = slide()
header(s, "07 · АДМИН", "Салбарын удирдлага", nextnum())
feature_card(s, 0.75, 1.95, 5.85, 3.0, "🏬", "Салбар бүртгэх",
             ["Шинэ салбар нэмж, нэр,", "утас, хаяг (хот/дүүрэг/", "хороо) оруулна.",
              "Газрын зураг дээр", "байршлыг тэмдэглэнэ."], accent=TEAL)
feature_card(s, 6.95, 1.95, 5.6, 3.0, "👁️", "Харуулах эсэх",
             ["«Салбаруудыг нийтэд", "харуулах» тохиргоогоор", "хэрэглэгчид харагдахыг",
              "удирдана. Салбар бүрийг", "идэвхтэй/идэвхгүй болгоно."], accent=VIOLET)
shot_slide("07 · АДМИН", "Салбарын тохиргоо",
           [("branch_settings_admin.png", "Салбар нэмэх/засах · «Салбаруудыг нийтэд харуулах» тохиргоо")])

# ════════════════════════════════════════════════════════════════════════════
# 12. ЗӨВЛӨМЖ / ТУСЛАМЖ
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "08 · ТУСЛАМЖ", "Зөвлөмж ба түгээмэл асуултууд", nextnum())
faqs = [
    ("Нэвтэрч чадахгүй байна?", "И-мэйл, нууц үгээ шалгаад «Нууц үг сэргээх»-ийг ашиглана."),
    ("Төлбөр төлсөн ч баталгаажихгүй?", "Хэдэн секунд хүлээгээд «Захиалгууд» хэсгээс төлөвөө шалгана."),
    ("Барааны нөөц харагдахгүй?", "Админ нөөцийн тоог зөв оруулсан эсэхийг шалгана."),
    ("Хүргэлтийн хаяг буруу?", "«Хаяг» хэсгээс засаж, газрын зураг дээр дахин тэмдэглэнэ."),
    ("Админ эрх авах?", "Системийн администратороор дамжуулан эрхээ тохируулна."),
    ("Тайлан гаргах?", "Админ → Санхүү → Тайлан, огноогоор шүүж Excel татна."),
]
x0, y0, cw, ch, gx, gy = 0.75, 1.95, 5.85, 1.45, 0.3, 0.18
for i, (q, a) in enumerate(faqs):
    x = x0 + (i % 2) * (cw + gx); y = y0 + (i // 2) * (ch + gy)
    box(s, x, y, cw, ch, fill=CARD, line=CARD2, line_w=1.0, radius=0.1)
    txt(s, x + 0.3, y + 0.18, cw - 0.6, 0.5, [[("❔ " + q, 13.5, TEAL_LT, True)]], space_after=0)
    txt(s, x + 0.3, y + 0.66, cw - 0.55, 0.7, [[(a, 11.5, GREY, False)]], space_after=0, line_spacing=1.05)

# ════════════════════════════════════════════════════════════════════════════
# 13. ТӨГСГӨЛ
# ════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "", "", nextnum())  # keep numbering; redraw clean
box(s, 9.3, -1.2, 5.5, 5.5, fill=CARD, radius=0.5)
box(s, 10.5, 4.3, 4.5, 4.5, fill=CARD, radius=0.5)
txt(s, 0.9, 2.4, 10, 0.5, [[("АМЖИЛТ ХҮСЬЕ!", 15, TEAL, True)]], space_after=0)
txt(s, 0.85, 2.95, 11, 1.2, [[("IShop-оо хялбар ашиглаарай", 44, WHITE, True)]], space_after=0)
txt(s, 0.9, 4.3, 10.5, 1.0,
    [[("Энэ гарын авлага хэрэглэгч болон админ хэсгийн үндсэн үйлдлүүдийг хамарлаа.", 16, GREY, False)],
     [("Нэмэлт асуулт байвал системийн администратортой холбогдоно уу.", 16, GREY, False)]],
    space_after=4, line_spacing=1.15)
chip(s, 0.9, 5.9, "IShop · 2026", TEAL)

# ── Хадгалах ────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(SHOT_DIR), "IShop_Хэрэглэгчийн_гарын_авлага.pptx")
prs.save(out)
with open(os.path.join(os.path.dirname(SHOT_DIR), "_guide_build.log"), "w", encoding="utf-8") as f:
    f.write(f"slides={len(prs.slides._sldIdLst)} out={out}\n")
